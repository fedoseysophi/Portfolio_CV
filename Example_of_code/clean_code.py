"""
Extract plain text from common office and PDF formats.
"""

from __future__ import annotations

import asyncio
import io
import os
import threading
from pathlib import Path
from typing import Callable, TypeVar

import aiofiles
import fitz
from docx import Document
from loguru import logger
from PIL import Image
import pandas as pd
import pytesseract
from pptx import Presentation

from app.config.schema import OCRConfig

T = TypeVar("T")

_OCR_ENV_NAMES = (
    "OMP_THREAD_LIMIT",
    "OPENBLAS_NUM_THREADS",
    "MKL_NUM_THREADS",
    "NUMEXPR_NUM_THREADS",
)
_tesseract_subprocess_env_lock = threading.Lock()


def _tesseract_image_to_string(
    img: Image.Image, lang: str, ocr_config: OCRConfig
) -> str:
    """
    Run Tesseract with capped BLAS/OpenMP threads and optional nice(1).

    pytesseract spawns a subprocess that inherits os.environ; concurrent OCR calls
    must not mutate env without a lock. High CPU usage during OCR is expected;
    thread_limit reduces multi-core spikes from one page; tesseract_nice yields CPU
    to other processes when set > 0 (Linux).
    """
    limit = str(ocr_config.thread_limit)
    with _tesseract_subprocess_env_lock:
        saved = {name: os.environ.get(name) for name in _OCR_ENV_NAMES}
        try:
            for name in _OCR_ENV_NAMES:
                os.environ[name] = limit
            return pytesseract.image_to_string(
                img,
                lang=lang,
                nice=ocr_config.tesseract_nice,
            )
        finally:
            for name in _OCR_ENV_NAMES:
                prev = saved[name]
                if prev is None:
                    os.environ.pop(name, None)
                else:
                    os.environ[name] = prev


async def _run_cpu_bound(func: Callable[..., T], *args: object, **kwargs: object) -> T:
    """Run a blocking function in the default thread pool."""
    return await asyncio.to_thread(func, *args, **kwargs)


def _extract_pdf_text_sync(file_path: str, ocr_config: OCRConfig) -> str:
    """Sync PDF extraction with optional per-page OCR."""
    parts: list[str] = []
    with fitz.open(file_path) as doc:
        for page_num, page in enumerate(doc):
            text = page.get_text()
            blocks = page.get_text("blocks")
            has_text_layer = any(b[6] == 0 and str(b[4]).strip() for b in blocks)
            skip_ocr = (
                not ocr_config.enabled
                or has_text_layer
                or len(text.strip()) >= ocr_config.min_chars
            )
            if skip_ocr:
                parts.append(f"--- Page {page_num + 1} ---\n{text}")
                continue
            logger.debug("Page {} has no text layer — running OCR", page_num + 1)
            pix = page.get_pixmap(dpi=ocr_config.dpi)
            img_bytes = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_bytes))
            ocr_text = _tesseract_image_to_string(
                img, lang=ocr_config.lang.replace("-", "_"), ocr_config=ocr_config
            )
            parts.append(f"--- Page {page_num + 1} (OCR) ---\n{ocr_text}")
    return "\n\n".join(parts)


def _extract_docx_text_sync(file_path: str) -> str:
    document = Document(file_path)
    return "\n".join(p.text for p in document.paragraphs if p.text.strip())


def _extract_pptx_text_sync(file_path: str) -> str:
    prs = Presentation(file_path)
    texts: list[str] = []
    for i, slide in enumerate(prs.slides):
        slide_text = [
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        ]
        texts.append(f"--- Slide {i + 1} ---\n{' '.join(slide_text)}")
    return "\n\n".join(texts)


def _extract_xlsx_text_sync(file_path: str) -> str:
    excel_file = pd.ExcelFile(file_path)
    blocks: list[str] = []
    for sheet_name in excel_file.sheet_names:
        df = excel_file.parse(sheet_name)
        blocks.append(f"--- Sheet: {sheet_name} ---\n{df.to_string()}")
    return "\n\n".join(blocks)


async def extract_pdf_text(file_path: str, ocr_config: OCRConfig) -> str:
    """Extract text from PDF; sparse pages may use OCR when enabled."""
    return await _run_cpu_bound(_extract_pdf_text_sync, file_path, ocr_config)


async def extract_docx_text(file_path: str) -> str:
    """Extract text from a Word document."""
    return await _run_cpu_bound(_extract_docx_text_sync, file_path)


async def extract_pptx_text(file_path: str) -> str:
    """Extract text from a PowerPoint file."""
    return await _run_cpu_bound(_extract_pptx_text_sync, file_path)


async def extract_xlsx_text(file_path: str) -> str:
    """Extract tabular content from Excel as plain text."""
    return await _run_cpu_bound(_extract_xlsx_text_sync, file_path)


async def extract_text_file(file_path: str) -> str:
    """Read a UTF-8 text file (best-effort decoding)."""
    async with aiofiles.open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return await f.read()


async def extract_text(file_path: str, ocr_config: OCRConfig) -> str:
    """
    Dispatch extraction based on file suffix.

    Args:
        file_path: Local filesystem path.
        ocr_config: OCR settings for PDF.

    Returns:
        Extracted plain text or empty string if unsupported / failed.
    """
    ext = Path(file_path).suffix.lower()
    try:
        if ext == ".pdf":
            return await extract_pdf_text(file_path, ocr_config)
        if ext == ".docx":
            return await extract_docx_text(file_path)
        if ext == ".pptx":
            return await extract_pptx_text(file_path)
        if ext == ".xlsx":
            return await extract_xlsx_text(file_path)
        if ext in {".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm"}:
            return await extract_text_file(file_path)
        logger.warning("No extractor for extension {} ({})", ext, file_path)
        return ""
    except Exception as e:
        logger.exception("Extraction failed for {}: {}", file_path, e)
        return ""
