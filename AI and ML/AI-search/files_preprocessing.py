import os
import re
import json
import spacy
import hashlib
from typing import List, Dict, Optional, Tuple
from tqdm import tqdm
from docx import Document
from pptx import Presentation
import fitz
import pandas as pd
import asyncio
import aiofiles
from concurrent.futures import ThreadPoolExecutor
import nest_asyncio
from datetime import datetime, timezone
import mimetypes
import time

# Enable event loop reentrancy (for Colab or similar)
nest_asyncio.apply()

# Load spaCy language model for sentence segmentation and lemmatization
global_nlp = spacy.load('ru_core_news_sm')
global_nlp.add_pipe('sentencizer')

class TextProcessor:
    def __init__(self, token_limit: int = 512):
        self.token_limit = token_limit

    @staticmethod
    def clean_text(text: str) -> str:
        # Remove unwanted characters from text
        text = re.sub(r'[^А-Яа-яЁё\s.,!?-]', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        return text

    def process_text(self, text: str, metadata: Dict) -> List[Dict]:
        cleaned_text = self.clean_text(text)
        doc = global_nlp(cleaned_text)
        sentences = list(doc.sents)

        chunks = []
        current_chunk = []
        current_token_count = 0
        chunk_number = 1

        # Split text into chunks based on token limit
        for sent in sentences:
            sent_token_count = len(sent)

            if current_token_count + sent_token_count > self.token_limit and current_chunk:
                chunk_text = " ".join([s.text for s in current_chunk])
                chunk_data = self._process_chunk(chunk_text, chunk_number, metadata)
                chunks.append(chunk_data)

                current_chunk = [sent]
                current_token_count = sent_token_count
                chunk_number += 1
            else:
                current_chunk.append(sent)
                current_token_count += sent_token_count

            # If single sentence exceeds token limit, store as separate chunk
            if sent_token_count > self.token_limit and len(current_chunk) == 1:
                chunk_text = current_chunk[0].text
                chunk_data = self._process_chunk(chunk_text, chunk_number, metadata)
                chunks.append(chunk_data)

                current_chunk = []
                current_token_count = 0
                chunk_number += 1

        # Add final chunk if not empty
        if current_chunk:
            chunk_text = " ".join([s.text for s in current_chunk])
            chunk_data = self._process_chunk(chunk_text, chunk_number, metadata)
            chunks.append(chunk_data)

        return chunks

    def _process_chunk(self, text: str, chunk_number: int, metadata: Dict) -> Dict:
        # Process individual chunk: extract entities, keyphrases, topics, stats, etc.
        doc = global_nlp(text)

        key_phrases = [
            token.text for token in doc
            if token.pos_ in ['NOUN', 'PROPN']
            and len(token.text.strip()) > 1
            and not token.is_stop
        ]
        lemmas = [
            token.lemma_ for token in doc
            if token.is_alpha and len(token.lemma_.strip()) > 1 and not token.is_stop
        ]
        entities = [
            {
                "text": ent.text,
                "label": ent.label_,
                "start": ent.start_char,
                "end": ent.end_char
            }
            for ent in doc.ents
        ]
        topics = [
            token.text for token in doc
            if token.pos_ in ['NOUN', 'PROPN'] and not token.is_stop
        ][:5]

        word_stats = self._generate_stats(doc)

        return {
            'chunk_number': chunk_number,
            'text': text,
            'token_count': len(doc),
            'entities': entities,
            'key_phrases': key_phrases,
            'lemmas': lemmas,
            'topics': topics,
            'stats': word_stats,
            'metadata': metadata,
            'processed_at': datetime.now().isoformat()
        }

    def _generate_stats(self, doc):
        # Count most frequent words in a chunk
        word_freq = {}
        for token in doc:
            if token.is_alpha and not token.is_stop:
                word_freq[token.lemma_] = word_freq.get(token.lemma_, 0) + 1
        return sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:5]

    def process_video_title(self, title: str, metadata: Dict) -> Dict:
        # Custom processing for video file titles
        cleaned_title = self.clean_text(title)
        doc = global_nlp(cleaned_title)

        lemmas = [
            token.lemma_ for token in doc
            if token.is_alpha and len(token.lemma_.strip()) > 1 and not token.is_stop
        ]
        key_phrases = [
            token.text for token in doc
            if token.pos_ in ['NOUN', 'PROPN'] and not token.is_stop
        ]
        topics = [
            token.text for token in doc
            if token.pos_ in ['NOUN', 'PROPN'] and not token.is_stop
        ][:5]

        return {
            'token_count': len(doc),
            'lemmas': lemmas,
            'key_phrases': key_phrases,
            'topics': topics,
            'metadata': metadata,
            'processed_at': datetime.now().isoformat()
        }

def calculate_hash(file_path: str) -> str:
    # Compute file hash (SHA256)
    hasher = hashlib.sha256()
    with open(file_path, 'rb') as f:
        hasher.update(f.read())
    return hasher.hexdigest()

def get_file_metadata(file_path: str) -> Dict:
    """Extract file metadata and infer semantic category."""
    try:
        file_stat = os.stat(file_path)
        file_name = os.path.basename(file_path)
        file_ext = os.path.splitext(file_name)[1].lower()

        # Detect MIME type
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type:
            # Basic fallback for common office extensions
            if file_ext in ['.docx']:
                mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif file_ext in ['.pptx']:
                mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            elif file_ext in ['.xlsx']:
                mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            elif file_ext in ['.pdf']:
                mime_type = 'application/pdf'
            elif file_ext in ['.mp4', '.avi', '.mov', '.mkv', '.webm']:
                mime_type = f'video/{file_ext[1:]}'
            else:
                mime_type = 'application/octet-stream'

        # Infer semantic category
        content_category = 'unknown'
        if mime_type:
            if mime_type.startswith('text/'):
                content_category = 'text'
            elif mime_type.startswith('image/'):
                content_category = 'image'
            elif mime_type.startswith('video/'):
                content_category = 'video'
            elif mime_type.startswith('audio/'):
                content_category = 'audio'
            elif 'pdf' in mime_type:
                content_category = 'document'
            elif 'word' in mime_type or 'document' in mime_type:
                content_category = 'document'
            elif 'presentation' in mime_type:
                content_category = 'presentation'
            elif 'sheet' in mime_type or 'excel' in mime_type:
                content_category = 'spreadsheet'

        return {
            'file_name': file_name,
            'file_path': file_path,
            'file_size': file_stat.st_size,
            'file_extension': file_ext,
            'mime_type': mime_type,
            'content_category': content_category,
            'created_at': datetime.fromtimestamp(file_stat.st_ctime).isoformat(),
            'modified_at': datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
            'accessed_at': datetime.fromtimestamp(file_stat.st_atime).isoformat(),
            'parent_directory': os.path.dirname(file_path)
        }
    except Exception as e:
        return {
            'file_name': os.path.basename(file_path),
            'file_path': file_path,
            'error': str(e)
        }

class DocumentProcessor:
    def __init__(self, directory: str, output_dir: str, token_limit: int = 512):
        self.directory = directory
        self.output_dir = output_dir
        self.text_processor = TextProcessor(token_limit)
        self.executor = ThreadPoolExecutor(max_workers=4)
        self.failed_files = []
        self.all_data = []
        self.log_file_path = os.path.join(self.output_dir, "processing_log.txt")

    async def process_files(self):
        """Main processing entrypoint for batch ingest."""
        self._write_log("Started file processing")
        start_time = datetime.now()
        self._write_log(f"Start time: {start_time.strftime('%Y-%m-%d %H:%M:%S')}")

        all_files = [
            os.path.join(root, file)
            for root, _, files in os.walk(self.directory)
            for file in files
        ]
        self._write_log(f"Total files found: {len(all_files)}")

        tasks = [self._process_file(file_path) for file_path in all_files]
        results = await asyncio.gather(*tasks)

        for res in results:
            if res is not None:
                self.all_data.append(res)

        await self._save_final_json()
        await self._save_final_csv()
        await self._save_failed_files()

        end_time = datetime.now()
        duration = end_time - start_time
        self._write_log(f"Processing finished: {end_time.strftime('%Y-%m-%d %H:%M:%S')}")
        self._write_log(f"Elapsed time: {str(duration)}")

    def _write_log(self, message: str):
        # Simple logger using plain text append
        os.makedirs(self.output_dir, exist_ok=True)
        with open(self.log_file_path, 'a', encoding='utf-8') as log_file:
            log_file.write(f"{datetime.now().strftime('%Y-%m-%d %H:%M:%S')} - {message}\n")

    async def _process_file(self, file_path: str) -> Optional[dict]:
        """Process a single file of any supported type."""
        try:
            current_hash = calculate_hash(file_path)
            if await self._is_cached(file_path, current_hash):
                return None

            file_metadata = get_file_metadata(file_path)
            file_name = file_metadata['file_name']
            file_name_without_ext = os.path.splitext(file_name)[0]
            ext = file_metadata['file_extension']

            processors = {
                '.pdf': self._process_pdf,
                '.docx': self._process_docx,
                '.pptx': self._process_pptx,
                '.xlsx': self._process_xlsx
            }

            if ext in processors:
                result = await processors[ext](file_path, file_metadata)
                await self._save_to_cache(file_path, result, current_hash)
                return {
                    "file_path": file_path,
                    "file_name": file_name_without_ext,
                    "metadata": file_metadata,
                    "chunks": result
                }
            elif ext in ['.mp4', '.avi', '.mov', '.mkv', '.webm']:
                video_analysis = self.text_processor.process_video_title(file_name_without_ext, file_metadata)
                video_chunk = {
                    'chunk_number': 1,
                    'text': file_name_without_ext,
                    'token_count': video_analysis.get('token_count', 0),
                    'entities': [],
                    'key_phrases': video_analysis["key_phrases"],
                    'lemmas': video_analysis["lemmas"],
                    'topics': video_analysis["topics"],
                    'stats': [],
                    'processed_at': datetime.now().isoformat()
                }
                result = [video_chunk]
                await self._save_to_cache(file_path, result, current_hash)
                return {
                    "file_path": file_path,
                    "file_name": file_name_without_ext,
                    "metadata": file_metadata,
                    "chunks": result
                }
            else:
                result = await self._process_generic_text(file_path, file_metadata)
                if result:
                    await self._save_to_cache(file_path, result, current_hash)
                    return {
                        "file_path": file_path,
                        "file_name": file_name_without_ext,
                        "metadata": file_metadata,
                        "chunks": result
                    }

        except FileNotFoundError:
            self.failed_files.append(file_path)
            self._write_log(f"File not found: {file_path}")
        except (PermissionError, IOError):
            self.failed_files.append(file_path)
            self._write_log(f"Access error: {file_path}")
        except Exception as e:
            self.failed_files.append(file_path)
            self._write_log(f"Failed to process: {file_path} - {str(e)}")
        return None

    async def _process_pdf(self, file_path: str, metadata: Dict) -> List[Dict]:
        # Parse PDF: extract text, metadata and split into segments
        text = ""
        with fitz.open(file_path) as pdf:
            page_count = len(pdf)
            metadata['page_count'] = page_count

            pdf_metadata = pdf.metadata
            if pdf_metadata:
                for key, value in pdf_metadata.items():
                    if value and isinstance(value, str):
                        metadata[f'pdf_{key.lower()}'] = value

            for page_num, page in enumerate(pdf):
                page_text = page.get_text()
                text += f"\n--- Page {page_num + 1} ---\n{page_text}"

        return self.text_processor.process_text(text, metadata)

    async def _process_docx(self, file_path: str, metadata: Dict) -> List[Dict]:
        # Parse Word document
        doc = Document(file_path)
        metadata['paragraph_count'] = len(doc.paragraphs)

        core_properties = doc.core_properties
        if core_properties:
            if core_properties.author:
                metadata['docx_author'] = core_properties.author
            if core_properties.title:
                metadata['docx_title'] = core_properties.title
            if core_properties.subject:
                metadata['docx_subject'] = core_properties.subject
            if core_properties.keywords:
                metadata['docx_keywords'] = core_properties.keywords
            if core_properties.created:
                metadata['docx_created'] = core_properties.created.isoformat() if hasattr(core_properties.created, 'isoformat') else str(core_properties.created)
            if core_properties.modified:
                metadata['docx_modified'] = core_properties.modified.isoformat() if hasattr(core_properties.modified, 'isoformat') else str(core_properties.modified)

        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        return self.text_processor.process_text(text, metadata)

    async def _process_pptx(self, file_path: str, metadata: Dict) -> List[Dict]:
        # Parse PowerPoint presentation
        prs = Presentation(file_path)
        metadata['slide_count'] = len(prs.slides)

        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                slides_text.append(f"--- Slide {i + 1} ---\n{' '.join(slide_text)}")
        text = "\n\n".join(slides_text)
        return self.text_processor.process_text(text, metadata)

    async def _process_xlsx(self, file_path: str, metadata: Dict) -> List[Dict]:
        # Parse Excel file
        excel_file = pd.ExcelFile(file_path)
        metadata['sheet_count'] = len(excel_file.sheet_names)
        metadata['sheet_names'] = excel_file.sheet_names

        text_all_sheets = []
        for sheet_name in excel_file.sheet_names:
            df = excel_file.parse(sheet_name=sheet_name)
            sheet_info = f"--- Sheet: {sheet_name} ---\nRows: {len(df)}, Columns: {len(df.columns)}\n"
            text_sheet = df.to_string(index=False, header=True)
            text_all_sheets.append(f"{sheet_info}\n{text_sheet}")

        full_text = "\n\n".join(text_all_sheets)
        return self.text_processor.process_text(full_text, metadata)

    async def _process_generic_text(self, file_path: str, metadata: Dict) -> Optional[List[Dict]]:
        # Parse generic text files
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                text = await f.read()
            metadata['line_count'] = text.count('\n') + 1
            metadata['char_count'] = len(text)
            metadata['word_count'] = len(text.split())
            return self.text_processor.process_text(text, metadata)
        except Exception as e:
            self._write_log(f"Failed to process text file: {file_path} - {str(e)}")
            return None

    async def _is_cached(self, file_path: str, current_hash: str) -> bool:
        # Check if file is already processed and cached
        base_filename = os.path.basename(file_path)
        cache_files = [
            f for f in os.listdir(self.output_dir)
            if f.startswith(base_filename) and f.endswith('.json')
        ]
        for cache_file in cache_files:
            if current_hash not in cache_file:
                await self._remove_old_hashes(os.path.join(self.output_dir, cache_file))

        cache_file = os.path.join(self.output_dir, f"{base_filename}_{current_hash}.json")
        return os.path.exists(cache_file)

    async def _remove_old_hashes(self, old_cache_file: str):
        # Remove outdated cache files
        try:
            os.remove(old_cache_file)
            self._write_log(f"Removed old cache: {old_cache_file}")
        except Exception as e:
            self._write_log(f"Failed to remove cache: {old_cache_file}")

    async def _save_to_cache(self, file_path: str, data: List[Dict], file_hash: str):
        # Cache the result of processing a file
        cache_file = os.path.join(self.output_dir, f"{os.path.basename(file_path)}_{file_hash}.json")
        os.makedirs(self.output_dir, exist_ok=True)
        async with aiofiles.open(cache_file, 'w', encoding='utf-8') as f:
            await f.write(json.dumps(data, ensure_ascii=False, indent=2))

    async def _save_final_json(self):
        # Save all results as a single JSON file
        final_json_path = os.path.join(self.output_dir, "all_files_result.json")
        async with aiofiles.open(final_json_path, 'w', encoding='utf-8') as f:
            await f.write(json.dumps(self.all_data, ensure_ascii=False, indent=2))

    async def _save_final_csv(self):
        # Export all results as a CSV
        csv_path = os.path.join(self.output_dir, "all_files_result.csv")
        csv_data = []

        for file_data in self.all_data:
            file_path = file_data["file_path"]
            file_name = file_data["file_name"]
            metadata = file_data.get("metadata", {})
            for chunk in file_data["chunks"]:
                row = {
                    "file_path": file_path,
                    "file_name": file_name,
                    "chunk_number": chunk["chunk_number"],
                    "token_count": chunk.get("token_count", 0),
                    "text": chunk["text"],
                    "processed_at": chunk.get("processed_at", "")
                }
                row["entities"] = "|".join([e["text"] for e in chunk.get("entities", [])])
                row["entity_types"] = "|".join([e["label"] for e in chunk.get("entities", [])])
                row["key_phrases"] = "|".join(chunk.get("key_phrases", []))
                row["lemmas"] = "|".join(chunk.get("lemmas", []))
                row["topics"] = "|".join(chunk.get("topics", []))
                row["stats"] = "|".join([f"{word}:{count}" for word, count in chunk.get("stats", [])])

                row["file_type"] = metadata.get("file_extension", "")
                row["content_category"] = metadata.get("content_category", "")
                row["mime_type"] = metadata.get("mime_type", "")
                row["file_size"] = metadata.get("file_size", 0)
                row["created_at"] = metadata.get("created_at", "")
                row["modified_at"] = metadata.get("modified_at", "")

                if "page_count" in metadata:
                    row["page_count"] = metadata["page_count"]
                if "docx_title" in metadata:
                    row["document_title"] = metadata["docx_title"]
                if "docx_author" in metadata:
                    row["document_author"] = metadata["docx_author"]
                if "slide_count" in metadata:
                    row["slide_count"] = metadata["slide_count"]
                if "sheet_count" in metadata:
                    row["sheet_count"] = metadata["sheet_count"]

                csv_data.append(row)

        df = pd.DataFrame(csv_data)
        df.to_csv(csv_path, index=False, encoding='utf-8')

    async def _save_failed_files(self):
        # Save failed files to log
        if self.failed_files:
            log_file = os.path.join(self.output_dir, 'failed_files.log')
            async with aiofiles.open(log_file, 'w', encoding='utf-8') as f:
                await f.write("\n".join(self.failed_files))

    async def create_search_index(self):
        """Creates a ready-to-use search index for the processed documents"""
        search_index_path = os.path.join(self.output_dir, "search_index.json")
        search_index = []
        for file_data in self.all_data:
            file_path = file_data["file_path"]
            file_name = file_data["file_name"]
            metadata = file_data.get("metadata", {})
            for chunk in file_data["chunks"]:
                search_entry = {
                    "id": f"{hashlib.md5(file_path.encode()).hexdigest()}_{chunk['chunk_number']}",
                    "file_path": file_path,
                    "file_name": file_name,
                    "chunk_number": chunk["chunk_number"],
                    "text": chunk["text"],
                    "key_phrases": chunk.get("key_phrases", []),
                    "lemmas": chunk.get("lemmas", []),
                    "topics": chunk.get("topics", []),
                    "content_category": metadata.get("content_category", ""),
                    "file_type": metadata.get("file_extension", ""),
                    "created_at": metadata.get("created_at", ""),
                    "modified_at": metadata.get("modified_at", "")
                }
                search_index.append(search_entry)

        async with aiofiles.open(search_index_path, 'w', encoding='utf-8') as f:
            await f.write(json.dumps(search_index, ensure_ascii=False, indent=2))
        self._write_log(f"Created search index with {len(search_index)} entries.")

if __name__ == "__main__":
    # Replace these with your input/output directories as needed
    input_dir = "/path/to/input"
    output_dir = "/path/to/output"

    processor = DocumentProcessor(input_dir, output_dir, token_limit=512)

    async def main():
        await processor.process_files()
        await processor.create_search_index()

    asyncio.run(main())
